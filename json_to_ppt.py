import json
from typing import List, Dict, Any, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re

class LyricsPPTGenerator:
    def __init__(self, json_path: str):
        self.json_path = json_path
        self.data = self._load_json()
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # 配置常量
        self.FONT_CN = "宋体"
        self.FONT_EN = "Times New Roman"
        self.COLOR_PRIMARY = RGBColor(0, 0, 0)  # 黑色
        self.COLOR_SECONDARY = RGBColor(80, 80, 80) # 深灰
        self.COLOR_STRESS = RGBColor(220, 20, 60)  # 深红色（重音）
        self.COLOR_SECONDARYSTRESS = RGBColor(240, 125, 150)  # 浅红色 （次重音）
        
        # 字号配置
        self.FONT_SIZE_JAPANESE = Pt(36)
        self.FONT_SIZE_ROMAJI = Pt(32)
        self.FONT_SIZE_TRANSLATION = Pt(28)
        self.FONT_SIZE_ANNOTATION = Pt(14)
    
    def _load_json(self) -> Dict[str, Any]:
        """加载 JSON 数据"""
        with open(self.json_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    
    def generate(self, output_path: str):
        """生成 PPT"""
        # 添加标题页
        self._add_title_slide()
        
        # 添加歌词页（跳过第一个段落：蝋燭/蜡烛）
        for verse in self.data['verses'][1:]:
            self._add_verse_slide(verse)
        
        # 保存
        self.prs.save(output_path)
        print(f"✅ PPT 生成完成: {output_path}")
    
    def _add_title_slide(self):
        """添加标题页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)
        
        # 日文标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        p = title_frame.paragraphs[0]
        p.text = self.data['metadata']['title']
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.name = self.FONT_CN
        p.alignment = PP_ALIGN.CENTER
        
        # 中文标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.0), Inches(9), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        
        p = subtitle_frame.paragraphs[0]
        p.text = self.data['metadata']['title_cn']
        p.font.size = Pt(44)
        p.font.name = self.FONT_CN
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = self.COLOR_SECONDARY
        
        # 作词作曲
        info_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.5), Inches(9), Inches(1.5)
        )
        info_frame = info_box.text_frame
        info_frame.word_wrap = True
        
        p = info_frame.paragraphs[0]
        p.text = f"作词: {self.data['metadata'].get('lyricist', '')}\n作曲: {self.data['metadata'].get('composer', '')}"
        p.font.size = Pt(20)
        p.font.name = self.FONT_CN
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = self.COLOR_SECONDARY
    
    def _add_verse_slide(self, verse: Dict[str, Any]):
        """添加歌词页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局
        
        # 背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # ========== 调整位置让罗马音位于中线 ==========
        # 日文歌词 (y: 1.2)
        japanese_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.0), Inches(9), Inches(0.8)
        )
        japanese_frame = japanese_box.text_frame
        japanese_frame.word_wrap = True
        
        p = japanese_frame.paragraphs[0]
        self._add_japanese_with_markers(p, verse)
        p.font.size = self.FONT_SIZE_JAPANESE
        p.font.bold = True
        p.font.name = self.FONT_CN
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)
        
        # 罗马音 (y: 2.1) - 中线位置
        romaji_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.0), Inches(9), Inches(0.8)
        )
        romaji_frame = romaji_box.text_frame
        romaji_frame.word_wrap = True
        
        p = romaji_frame.paragraphs[0]
        self._add_romaji_with_stress(p, verse['romaji'])
        p.font.size = self.FONT_SIZE_ROMAJI
        p.font.name = self.FONT_EN
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)
        
        # 翻译 (y: 3.0)
        translation_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.0), Inches(9), Inches(0.6)
        )
        translation_frame = translation_box.text_frame
        translation_frame.word_wrap = True
        
        p = translation_frame.paragraphs[0]
        p.text = verse['translation']
        p.font.size = self.FONT_SIZE_TRANSLATION
        p.font.name = self.FONT_CN
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = self.COLOR_SECONDARY
        p.space_after = Pt(4)
        
        # ========== 注释区域（更紧凑） ==========
        if verse['annotations']:
            annotations_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(5.0), Inches(8.6), Inches(3.4)
            )
            annotations_frame = annotations_box.text_frame
            annotations_frame.word_wrap = True
            
            # 清除默认段落
            annotations_frame.clear()
            
            for i, annotation in enumerate(verse['annotations']):
                # 添加注释段落
                p = annotations_frame.add_paragraph()
                self._add_annotation_paragraph(p, annotation)
                p.space_before = Pt(0)
                p.space_after = Pt(6)  # 减少间距
    
    def _add_japanese_with_markers(self, paragraph, verse: Dict[str, Any]):
        """
        添加日文文本，在注释词下加下划线，并在词后添加注释标记
        例: パッと①咲いた②
        """
        japanese = verse['japanese']
        annotations = verse['annotations']
        
        # 构建一个映射：词 -> 标记
        term_to_marker = {}
        for annotation in annotations:
            if annotation['term']:
                term_to_marker[annotation['term']] = annotation['marker']
        
        if not term_to_marker:
            # 如果没有注释词，直接添加日文
            paragraph.clear()
            run = paragraph.add_run()
            run.text = japanese
            run.font.name = self.FONT_CN
            return
        
        # 清空段落
        paragraph.clear()
        
        # 遍历日文，逐个处理
        i = 0
        while i < len(japanese):
            # 检查是否有匹配的注释词（优先匹配最长）
            matched = False
            for term in sorted(term_to_marker.keys(), key=len, reverse=True):
                if japanese[i:].startswith(term):
                    # 添加下划线的词
                    run = paragraph.add_run()
                    run.text = term
                    run.font.name = self.FONT_CN
                    run.font.underline = True
                    
                    # 添加标记（正常大小）
                    marker = term_to_marker[term]
                    run_marker = paragraph.add_run()
                    run_marker.text = marker
                    run_marker.font.name = self.FONT_CN
                    
                    i += len(term)
                    matched = True
                    break
            
            if not matched:
                # 普通字符
                run = paragraph.add_run()
                run.text = japanese[i]
                run.font.name = self.FONT_CN
                i += 1
    def _add_romaji_with_stress(self, paragraph, romaji_list: List[Dict[str, str]]):
        """
        添加罗马音，重音用粗体+颜色标记
        """
        paragraph.clear()
        
        for item in romaji_list:
            run = paragraph.add_run()
            run.text = item['text']
            run.font.name = self.FONT_EN
            
            if item['stress'] == 'primary':
                # 粗体 + 深红色
                run.font.bold = True
                run.font.color.rgb = self.COLOR_STRESS
            elif item['stress'] == 'secondary':
                # 粗体 + 浅红色
                run.font.bold = True
                run.font.color.rgb = self.COLOR_SECONDARYSTRESS
            else:
                # 普通文本
                run.font.color.rgb = self.COLOR_PRIMARY
    
    def _add_annotation_paragraph(self, paragraph, annotation: Dict[str, str]):
        """
        添加单条注释段落
        格式: ①term（pronunciation）：content
        """
        # 标记（不变色，保持黑色）
        run = paragraph.add_run()
        run.text = annotation['marker']
        run.font.size = self.FONT_SIZE_ANNOTATION
        run.font.bold = True
        run.font.name = self.FONT_CN
        run.font.color.rgb = self.COLOR_PRIMARY
        
        # 词（不下划线）
        if annotation['term']:
            run = paragraph.add_run()
            run.text = annotation['term']
            run.font.size = self.FONT_SIZE_ANNOTATION
            run.font.name = self.FONT_CN
            run.font.color.rgb = self.COLOR_PRIMARY
            
            # 发音（如果存在）
            if annotation['pronunciation']:
                run = paragraph.add_run()
                run.text = f"（{annotation['pronunciation']}）"
                run.font.size = Pt(13)
                run.font.name = self.FONT_EN
                run.font.color.rgb = self.COLOR_SECONDARY
            
            # 冒号
            run = paragraph.add_run()
            run.text = "："
            run.font.size = self.FONT_SIZE_ANNOTATION
            run.font.name = self.FONT_CN
            run.font.color.rgb = self.COLOR_PRIMARY
        
        # 内容
        run = paragraph.add_run()
        run.text = annotation['content']
        run.font.size = self.FONT_SIZE_ANNOTATION
        run.font.name = self.FONT_CN
        run.font.color.rgb = self.COLOR_PRIMARY
        
        # 段落格式
        paragraph.level = 0
        paragraph.line_spacing = 1.25


if __name__ == "__main__":
    json_path = r'c:\Users\lenovo\Desktop\日语学习\歌词_example.json'
    output_path = r'c:\Users\lenovo\Desktop\日语学习\歌词_example.pptx'
    
    try:
        generator = LyricsPPTGenerator(json_path)
        generator.generate(output_path)
        print(f"✅ 成功生成 PPT")
        print(f"📊 文件位置: {output_path}")
    except Exception as e:
        print(f"❌ 错误: {e}")
        import traceback
        traceback.print_exc()